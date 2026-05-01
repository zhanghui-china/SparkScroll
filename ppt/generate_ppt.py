#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
SparkScroll 项目 PPT 生成脚本
使用 python-pptx 库生成专业的演示文稿
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# 创建演示文稿对象
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# NVIDIA 绿色主题
NVIDIA_GREEN = RGBColor(118, 185, 0)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(240, 240, 240)

def add_title_slide(title, subtitle=""):
    """添加标题页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加绿色背景条
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(10), Inches(2.5)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = NVIDIA_GREEN
    shape.line.fill.background()
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(9), Inches(1))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 添加副标题
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(0.5))
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(title, content_items, has_icon=False):
    """添加内容页"""
    slide_layout = prs.slide_layouts[6]  # 空白布局
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NVIDIA_GREEN
    
    # 添加内容
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)
        p.level = 0
    
    return slide

def add_two_column_slide(title, left_title, left_items, right_title, right_items):
    """添加双栏内容页"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    # 添加标题
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = NVIDIA_GREEN
    
    # 左栏标题
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(4.2), Inches(0.5))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # 左栏内容
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(4.2), Inches(4.8))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    # 右栏标题
    right_title_box = slide.shapes.add_textbox(Inches(5.3), Inches(1.5), Inches(4.2), Inches(0.5))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_GRAY
    
    # 右栏内容
    right_content = slide.shapes.add_textbox(Inches(5.3), Inches(2.2), Inches(4.2), Inches(4.8))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)
    
    return slide

# ==================== 开始生成幻灯片 ====================

# 第1页：封面
slide1 = add_title_slide(
    "SparkScroll - 星火绘卷",
    "基于 NVIDIA DGX Spark 的智能连环画生成平台\n\n团队：纵贯线 | 时间：2026年4月"
)

# 第2页：项目目标
add_content_slide(
    "🎯 项目目标",
    [
        "让每个人都能成为漫画家",
        "用 AI 重铸中国连环画的数字辉煌",
        "",
        "核心价值：",
        "  ✅ 降低创作门槛：从\"绘画技法\"到\"叙事审美\"",
        "  ✅ 缩短制作周期：从数月到数小时",
        "  ✅ 压缩成本：从数百美元到接近零成本"
    ]
)

# 第3页：五大应用场景
add_two_column_slide(
    "💡 五大应用场景",
    "痛点",
    [
        "1️⃣ K-12 智慧教育",
        "   古典名著晦涩难懂",
        "",
        "2️⃣ 网文 IP 极速孵化",
        "   漫改周期长，成本高",
        "",
        "3️⃣ 数字文博策展",
        "   视觉资产匮乏",
        "",
        "4️⃣ 品牌出海营销",
        "   难以响应热点",
        "",
        "5️⃣ 连环画文化复兴",
        "   产业逐渐断层"
    ],
    "解决方案",
    [
        "将《西游记》转化为连环画",
        "知识留存率提升 40%",
        "",
        "每日更新自动生成连环画",
        "社媒快速引流",
        "",
        "县志转化为历史连环长卷",
        "沉浸式展览",
        "",
        "品牌故事快速转化",
        "多风格多语言适配",
        "",
        "AI 重建文化生产线",
        "\"小人书\"涅槃重生"
    ]
)

# 第4页：六阶段工作流 - 第1部分
add_content_slide(
    "⚙️ 六阶段工业级工作流 (1/2)",
    [
        "1️⃣ 导演 Agent (Director)",
        "    • 文本降维压缩，提取核心剧情和人物",
        "    • 基于 Qwen3.5-9B，128K 上下文",
        "",
        "2️⃣ 编剧 Agent (Writer)",
        "    • 分镜拆解，每页 4-6 个场景",
        "    • 输出标准化 JSON 剧本",
        "",
        "3️⃣ 立绘 Agent (Character Designer) ⭐",
        "    • 预生成角色三视图和道具参考图",
        "    • 核心突破：解决角色一致性难题"
    ]
)

# 第5页：六阶段工作流 - 第2部分
add_content_slide(
    "⚙️ 六阶段工业级工作流 (2/2)",
    [
        "4️⃣ 出稿 Agent (Drafting)",
        "    • 基于分镜 + 角色参考图生成底图",
        "    • Diffusers + vLLM-Omni 双驱动",
        "",
        "5️⃣ 剪辑 Agent (Editor)",
        "    • Python PIL 物理渲染文字",
        "    • 拒绝 AI 生字乱码，印刷级清晰度",
        "",
        "6️⃣ 合成 Agent (Assembler)",
        "    • 流式按页呈现，即生成即预览",
        "    • 理想状态：导出连贯 PDF 画册"
    ]
)

# 第6页：角色一致性解决方案
add_content_slide(
    "🔑 角色一致性解决方案",
    [
        "传统问题：",
        "  ✗ AI 生成的角色在不同分镜中形象不连贯",
        "  ✗ 服饰、容貌、气质频繁变化",
        "  ✗ 用户阅读体验断裂",
        "",
        "我们的方案：",
        "  ✓ 立绘 Agent 预生成角色参考图",
        "      - 正面、侧面、背面三视图",
        "      - 关键道具标准参考图",
        "  ✓ 后续所有分镜以参考图为基准",
        "      - 注入扩散模型条件通道",
        "      - 风格对齐，形象统一",
        "",
        "效果：角色高度一致 | 叙事连贯流畅 | 专业出版级质量"
    ]
)

# 第7页：技术架构
add_content_slide(
    "🏗️ 技术架构",
    [
        "前端层：Streamlit UI",
        "  ├── 项目管理 | 进度监控 | 页面预览",
        "",
        "网关层：FastAPI Gateway",
        "  ├── API 接口 | 任务调度 | 状态管理",
        "",
        "模型层：",
        "  ├── 文本：Qwen3.5-9B (vLLM v0.19.0)",
        "  ├── 图像：Qwen-Image-Edit-2511 (Diffusers + vLLM-Omni)",
        "",
        "运行态：Redis",
        "  ├── 任务队列 | 锁 | 运行态数据",
        "",
        "存储层：本地文件系统",
        "  ├── 项目资源隔离"
    ]
)

# 第8页：DGX Spark 不可替代性
add_content_slide(
    "🚀 为什么必须用 NVIDIA DGX Spark？",
    [
        "重型双大模型常驻显存架构",
        "",
        "  • 文本主脑：Qwen3.5-9B (128K) ─────────── 40GB",
        "  • 视觉中枢：Qwen-Image-Edit-2511 ──────── 50-60GB",
        "  • 框架系统：vLLM + FastAPI + Diffusers ── 15GB",
        "  • 总计：110-115GB / 128GB",
        "",
        "普通显卡的困境：",
        "  ✗ RTX 4090 (24GB)：显存不足，频繁 Swap",
        "  ✗ PCIe 多卡：总线传输瓶颈，延迟断崖式下跌",
        "  ✗ 长上下文：速度衰减可达 64 倍",
        "",
        "DGX Spark 的优势：",
        "  ✓ 128GB 统一内存 (Grace CPU + Blackwell GPU)",
        "  ✓ NVLink-C2C：900 GB/s 双向带宽",
        "  ✓ 零 Swap 极致响应 | 多 Agent 极速切换",
        "",
        "结论：普通单卡系统完全无法运行该并发管线！"
    ]
)

# 第9页：NVIDIA 技术栈
add_two_column_slide(
    "🛠️ 使用的 NVIDIA 技术",
    "平台与硬件",
    [
        "• NVIDIA DGX Spark (GB10)",
        "• 128GB LPDDR5x 统一内存",
        "• Grace CPU + Blackwell GPU",
        "• NVLink-C2C 900 GB/s"
    ],
    "软件栈",
    [
        "• CUDA 12.9/13.0",
        "• PyTorch 2.10",
        "• vLLM v0.19.0",
        "  - PagedAttention",
        "  - 显存利用率 40%→90%",
        "  - 吞吐量提升 2-4 倍",
        "• vLLM-Omni",
        "• Diffusers",
        "",
        "模型：",
        "• Qwen3.5-9B (128K)",
        "• Qwen-Image-Edit-2511",
        "• FireRed-Image-Edit-1.1"
    ]
)

# 第10页：演示效果展示
slide10 = add_title_slide(
    "🎨 演示效果展示",
    "案例1：《卖火柴的小女孩》 - 20页精美连环画\n案例2：《小红帽》 - 连贯多页连环画\n\n[此处插入连环画样张图片]"
)

# 第11页：项目成果
add_content_slide(
    "📊 项目成果",
    [
        "✅ 六阶段工作流完整实现",
        "    导演 → 编剧 → 立绘 → 出稿 → 剪辑 → 合成",
        "",
        "✅ 双本地模型部署",
        "    Qwen3.5-9B (文本) + Qwen-Image-Edit-2511 (图像)",
        "",
        "✅ 前端界面开发",
        "    Streamlit 流式预览 | 项目管理 | 进度监控",
        "",
        "✅ API 网关与调度系统",
        "    FastAPI 接口 | Redis 运行态管理 | 任务队列",
        "",
        "✅ 角色一致性方案验证",
        "    立绘预生成 | 参考图注入"
    ]
)

# 第12页：未来规划
add_two_column_slide(
    "🚀 未来规划",
    "短期（当前）",
    [
        "• 继续完善剪辑 Agent",
        "• 完善合成 Agent",
        "• 漫画风格模型训练",
        "• 增强分镜连续性"
    ],
    "长期",
    [
        "• 模型扩展",
        "  - 视频生成",
        "  - 3D 内容创作",
        "",
        "• 性能优化",
        "  - 模型量化、剪枝",
        "",
        "• 云服务版本",
        "• 行业应用拓展",
        "  - 广告、教育、游戏",
        "",
        "• 开源社区建设"
    ]
)

# 第13页：团队介绍
add_content_slide(
    "👥 纵贯线团队",
    [
        "团队成员：",
        "  • 张辉 (张小白) - 队长、项目策划",
        "      场景策划 | 模型部署 | 前端开发",
        "",
        "  • 寒晨 - 架构设计、代码负责",
        "      系统架构 | Gateway 开发",
        "",
        "  • 覃飞雄 (飞哥) - 测试与文档",
        "      前端开发 | 系统测试 | 文档编写",
        "",
        "  • Codex - Gateway 开发",
        "  • Trae - 前端开发",
        "",
        "协作亮点：",
        "  ✓ 跨越北京、南京、广州三地协作",
        "  ✓ OpenClaw 多 Agent 环境辅助开发",
        "  ✓ 远程共享 DGX Spark 设备"
    ]
)

# 第14页：结语
slide14 = add_title_slide(
    "🌟 星火绘卷的意义",
    "不仅是一次算力破局\n更是一次横跨时代的文化接力\n\n📚 古典文学遗产 → 数字画卷\n🎨 AI 技术 → 文化复兴\n🚀 工业级生产 → 零门槛创作\n\n中国连环画，正以数字化的崭新形态涅槃重生！"
)

# 第15页：联系方式
add_content_slide(
    "📮 联系我们",
    [
        "项目代码仓库：",
        "https://gitee.com/zhanghui_china/SparkScroll",
        "",
        "云端演示环境：",
        "https://sparkscroll.opsx.info",
        "",
        "团队：纵贯线",
        "",
        "",
        "感谢：",
        "  • NVIDIA DGX Spark 平台提供的强大算力支持",
        "  • 黑客松组委会提供的展示机会",
        "  • 团队成员的辛勤付出"
    ]
)

# 保存文件
output_path = "./SparkScroll_演示文稿.pptx"
prs.save(output_path)
print(f"✅ PPT 生成成功！")
print(f"📄 文件路径：{output_path}")
